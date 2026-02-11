rom pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.dml import MSO_THEME_COLOR

class CorporatePresentation:
    """
    Enhanced corporate presentation template with professional slide types.
    Colors: Deep Teal (#2C5F7C), Bright Red (#E31E24), Light Blue (#7BA7BC)
    Style: Modern, clean, professional layouts for sales and business presentations
    Features: Slide numbers, Waldom branding, responsive layouts, multiple chart types
    """
    
    def __init__(self):
        self.prs = Presentation()
        
        # Set slide dimensions (16:9 widescreen)
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)
        
        # Brand colors
        self.TEAL = RGBColor(44, 95, 124)
        self.RED = RGBColor(227, 30, 36)
        self.LIGHT_BLUE = RGBColor(123, 167, 188)
        self.WHITE = RGBColor(255, 255, 255)
        self.LIGHT_GRAY = RGBColor(242, 242, 242)
        self.DARK_GRAY = RGBColor(51, 51, 51)
        self.DARK_TEAL = RGBColor(30, 70, 95)
        
        self.FONT_NAME = "Calibri"
    
    def _add_footer(self, slide, slide_number):
        """Adds the slide number to the bottom right."""
        left = Inches(8.5)
        top = Inches(7.0)
        width = Inches(1.0)
        height = Inches(0.5)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        p = text_frame.paragraphs[0]
        p.text = str(slide_number)
        p.font.name = self.FONT_NAME
        p.font.size = Pt(10)
        p.font.color.rgb = self.DARK_GRAY
        p.alignment = PP_ALIGN.RIGHT

    # ========== TITLE & TABLE OF CONTENTS SLIDES ==========
    
    def add_title_slide(self, title, subtitle, slide_number=None):
        """Create title slide with Waldom logo and red accent, matching preview."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Background gradient effect (using shapes for similar visual)
        bg_shape_blue = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(10), Inches(7.5)
        )
        bg_shape_blue.fill.solid()
        bg_shape_blue.fill.fore_color.rgb = self.TEAL
        bg_shape_blue.line.fill.background()

        # Add "opacity" effect with lighter shapes
        oval1 = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(6), Inches(-0.5),
            Inches(4), Inches(4)
        )
        oval1.fill.solid()
        oval1.fill.fore_color.rgb = self.WHITE
        oval1.fill.fore_color.alpha = 0.05
        oval1.line.fill.background()

        oval2 = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(-0.5), Inches(5),
            Inches(3.5), Inches(3.5)
        )
        oval2.fill.solid()
        oval2.fill.fore_color.rgb = self.WHITE
        oval2.fill.fore_color.alpha = 0.05
        oval2.line.fill.background()

        # Waldom Logo (optional - will need actual logo path)
        logo_path = "path/to/your/waldom_logo.png"
        try:
            slide.shapes.add_picture(logo_path, Inches(3.8), Inches(1.8), height=Inches(0.8))
        except FileNotFoundError:
            print(f"Note: Waldom logo not found at {logo_path}. Continuing without logo.")

        # Title text
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(3),
            Inches(8), Inches(1.5)
        )
        tf = title_box.text_frame
        tf.text = title
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.font.name = self.FONT_NAME
        p.font.size = Pt(40)
        p.font.bold = False
        p.font.color.rgb = self.WHITE
        p.alignment = PP_ALIGN.CENTER

        # Subtitle
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(
                Inches(1), Inches(4.5),
                Inches(8), Inches(1)
            )
            stf = subtitle_box.text_frame
            stf.text = subtitle
            stf.word_wrap = True
            sp = stf.paragraphs[0]
            sp.font.name = self.FONT_NAME
            sp.font.size = Pt(16)
            sp.font.color.rgb = self.WHITE
            sp.alignment = PP_ALIGN.CENTER
        
        # Red line at bottom right
        red_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(8.5), Inches(6.8),
            Inches(1), Inches(0.05)
        )
        red_line.fill.solid()
        red_line.fill.fore_color.rgb = self.RED
        red_line.line.fill.background()

        if slide_number:
            self._add_footer(slide, slide_number)
        return slide
    
    def add_table_of_contents(self, sections, slide_number=None):
        """Create table of contents slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Header strip
        header_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(10), Inches(0.08)
        )
        header_line.fill.solid()
        header_line.fill.fore_color.rgb = self.TEAL
        header_line.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.8),
            Inches(5), Inches(0.8)
        )
        tf = title_box.text_frame
        tf.text = "Table Of Content"
        
        p = tf.paragraphs[0]
        p.font.name = self.FONT_NAME
        p.font.size = Pt(32)
        p.font.color.rgb = self.DARK_TEAL
        
        # Red accent line under title
        red_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5), Inches(1.6),
            Inches(4), Inches(0.05)
        )
        red_line.fill.solid()
        red_line.fill.fore_color.rgb = self.RED
        red_line.line.fill.background()
        
        # Add sections in two columns
        left_x = Inches(0.5)
        right_x = Inches(5)
        start_y = Inches(2.3)
        spacing = Inches(1.1)
        
        for i, section in enumerate(sections):
            if i % 2 == 0:
                x_pos = left_x
                y_pos = start_y + (i // 2) * spacing
            else:
                x_pos = right_x
                y_pos = start_y + (i // 2) * spacing
            
            # Number
            num_box = slide.shapes.add_textbox(x_pos, y_pos, Inches(0.8), Inches(0.5))
            ntf = num_box.text_frame
            ntf.text = f"{i+1:02d}."
            
            np = ntf.paragraphs[0]
            np.font.name = self.FONT_NAME
            np.font.size = Pt(28)
            np.font.bold = True
            np.font.color.rgb = self.DARK_GRAY
            
            # Section title
            text_box = slide.shapes.add_textbox(
                x_pos, y_pos + Inches(0.5),
                Inches(4), Inches(0.5)
            )
            ttf = text_box.text_frame
            ttf.text = section
            ttf.word_wrap = True
            
            tp = ttf.paragraphs[0]
            tp.font.name = self.FONT_NAME
            tp.font.size = Pt(14)
            tp.font.color.rgb = self.DARK_GRAY
        
        if slide_number:
            self._add_footer(slide, slide_number)
        return slide

    # ========== CONTENT SLIDES WITH ICONS ==========
    
    def add_content_with_icons_slide(self, title, items, slide_number=None):
        """Create content slide with icon-text pairs in responsive grid layout"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Header strip
        header_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(10), Inches(0.08)
        )
        header_line.fill.solid()
        header_line.fill.fore_color.rgb = self.TEAL
        header_line.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = self.FONT_NAME
        p.font.size = Pt(32)
        p.font.color.rgb = self.DARK_TEAL
        
        # Red line under title
        red_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5), Inches(1.6),
            Inches(4), Inches(0.05)
        )
        red_line.fill.solid()
        red_line.fill.fore_color.rgb = self.RED
        red_line.line.fill.background()

        # Content area starts lower
        start_y = Inches(2.2) 
        
        # Determine scaling based on item count
        item_count = len(items)
        icon_size_inch = 0.5 if item_count <= 4 else 0.4
        icon_font_size = Pt(24) if item_count <= 4 else Pt(20)
        text_font_size = Pt(18) if item_count <= 4 else Pt(16)
        line_spacing = 0.2 if item_count <= 4 else 0.15

        col_width = Inches(4.5)
        item_height = Inches(icon_size_inch + 0.5)
        
        for i, item in enumerate(items):
            col = i % 2
            row = i // 2
            
            x = Inches(0.5) + col * (col_width + Inches(0.5))
            y = start_y + row * (item_height + Inches(line_spacing))

            # Icon box
            icon_box = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                x, y,
                Inches(icon_size_inch), Inches(icon_size_inch)
            )
            icon_box.fill.solid()
            icon_box.fill.fore_color.rgb = self.RED
            icon_box.line.fill.background()

            # Icon text (placeholder)
            icon_text_box = slide.shapes.add_textbox(
                x, y,
                Inches(icon_size_inch), Inches(icon_size_inch)
            )
            itf = icon_text_box.text_frame
            itf.text = item.get('icon', '★')
            itf.paragraphs[0].font.name = self.FONT_NAME
            itf.paragraphs[0].font.size = icon_font_size
            itf.paragraphs[0].font.color.rgb = self.WHITE
            itf.paragraphs[0].alignment = PP_ALIGN.CENTER
            itf.vertical_anchor = MSO_ANCHOR.MIDDLE

            # Text
            text_box = slide.shapes.add_textbox(
                x + Inches(icon_size_inch + 0.2), y,
                col_width - Inches(icon_size_inch + 0.2), Inches(icon_size_inch + 0.5)
            )
            ttf = text_box.text_frame
            p = ttf.paragraphs[0]
            p.text = item['text']
            p.font.name = self.FONT_NAME
            p.font.size = text_font_size
            p.font.color.rgb = self.DARK_GRAY
            p.word_wrap = True
        
        if slide_number:
            self._add_footer(slide, slide_number)
        return slide

    # ========== TEXT & PARAGRAPH SLIDES ==========
    
    def add_split_slide(self, title, paragraphs, slide_number=None):
        """Create slide with title and multiple paragraphs"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Header strip
        header_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(10), Inches(0.08)
        )
        header_line.fill.solid()
        header_line.fill.fore_color.rgb = self.TEAL
        header_line.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = self.FONT_NAME
        p.font.size = Pt(32)
        p.font.color.rgb = self.DARK_TEAL
        
        # Red line under title
        red_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5), Inches(1.6),
            Inches(4), Inches(0.05)
        )
        red_line.fill.solid()
        red_line.fill.fore_color.rgb = self.RED
        red_line.line.fill.background()

        # Determine scaling based on paragraph count
        para_count = len(paragraphs)
        font_size = Pt(18) if para_count <= 2 else Pt(16) if para_count <= 3 else Pt(14)
        line_spacing = 0.2 if para_count <= 2 else 0.15 if para_count <= 3 else 0.1

        current_y = Inches(2.2)
        for para in paragraphs:
            text_box = slide.shapes.add_textbox(Inches(0.5), current_y, Inches(9), Inches(1))
            tf = text_box.text_frame
            p = tf.paragraphs[0]
            p.text = para
            p.font.name = self.FONT_NAME
            p.font.size = font_size
            p.font.color.rgb = self.DARK_GRAY
            p.word_wrap = True
            current_y += (text_box.height + Inches(line_spacing))
        
        if slide_number:
            self._add_footer(slide, slide_number)
        return slide

    # ========== MARKET OPPORTUNITIES & GRID LAYOUTS ==========
    
    def add_market_opportunities_slide(self, title, items, slide_number=None):
        """Create market opportunities slide with responsive grid"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Header strip
        header_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(10), Inches(0.08)
        )
        header_line.fill.solid()
        header_line.fill.fore_color.rgb = self.TEAL
        header_line.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = self.FONT_NAME
        p.font.size = Pt(32)
        p.font.color.rgb = self.DARK_TEAL
        
        # Red line under title
        red_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5), Inches(1.6),
            Inches(4), Inches(0.05)
        )
        red_line.fill.solid()
        red_line.fill.fore_color.rgb = self.RED
        red_line.line.fill.background()

        item_count = len(items)
        if item_count == 0:
            if slide_number:
                self._add_footer(slide, slide_number)
            return slide

        # Dynamic grid calculations
        grid_config = {
            1: (1, 1), 2: (1, 2), 3: (1, 3),
            4: (2, 2), 5: (2, 3), 6: (2, 3)
        }
        num_rows, max_cols_per_row = grid_config.get(item_count, (2, 3))
        
        # Adjust spacing and sizes
        if item_count <= 3:
            col_width = Inches(8.5 / item_count)
            item_h_padding = Inches(0.4)
            icon_size_inch = Inches(0.7)
            title_font_size = Pt(18)
            text_font_size = Pt(14)
            gap = Inches(0.4)
        else:
            col_width = Inches(8.5 / 3)
            item_h_padding = Inches(0.3)
            icon_size_inch = Inches(0.6)
            title_font_size = Pt(16)
            text_font_size = Pt(12)
            gap = Inches(0.3)
        
        start_x = Inches(0.5)
        start_y = Inches(2.2)
        current_item_idx = 0

        for r in range(num_rows):
            cols_in_this_row = max_cols_per_row
            if item_count == 5 and r == 1:
                cols_in_this_row = 2
            elif item_count <= 3:
                cols_in_this_row = item_count

            row_start_y = start_y + r * (Inches(2.2) + gap)

            for c in range(cols_in_this_row):
                if current_item_idx >= item_count:
                    break

                item = items[current_item_idx]
                
                x = start_x + c * (col_width + gap)
                y = row_start_y

                # Item box with border
                item_shape = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    x, y,
                    col_width, Inches(2.2)
                )
                item_shape.fill.background()
                item_shape.line.fill.solid()
                item_shape.line.fill.fore_color.rgb = self.LIGHT_BLUE
                item_shape.line.width = Pt(1.5)

                # Icon box
                icon_box_x = x + item_h_padding
                icon_box_y = y + item_h_padding
                icon_rect = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    icon_box_x, icon_box_y,
                    icon_size_inch, icon_size_inch
                )
                icon_rect.fill.solid()
                icon_rect.fill.fore_color.rgb = self.RED
                icon_rect.line.fill.background()

                # Icon placeholder
                icon_text_box = slide.shapes.add_textbox(
                    icon_box_x, icon_box_y,
                    icon_size_inch, icon_size_inch
                )
                itf = icon_text_box.text_frame
                itf.text = item.get('icon', '★')
                itf.paragraphs[0].font.name = self.FONT_NAME
                itf.paragraphs[0].font.size = Pt(14)
                itf.paragraphs[0].font.color.rgb = self.WHITE
                itf.paragraphs[0].alignment = PP_ALIGN.CENTER
                itf.vertical_anchor = MSO_ANCHOR.MIDDLE

                # Text content
                text_content_x = x + icon_size_inch + Inches(0.2) + item_h_padding
                text_content_width = col_width - icon_size_inch - Inches(0.2) - 2 * item_h_padding
                
                title_text_box = slide.shapes.add_textbox(
                    text_content_x, y + item_h_padding,
                    text_content_width, Inches(0.5)
                )
                ttf = title_text_box.text_frame
                p = ttf.paragraphs[0]
                p.text = item['title']
                p.font.name = self.FONT_NAME
                p.font.size = title_font_size
                p.font.color.rgb = self.DARK_TEAL
                p.word_wrap = True

                description_text_box = slide.shapes.add_textbox(
                    text_content_x, y + item_h_padding + Inches(0.5),
                    text_content_width, Inches(1)
                )
                dtf = description_text_box.text_frame
                p = dtf.paragraphs[0]
                p.text = item['text']
                p.font.name = self.FONT_NAME
                p.font.size = text_font_size
                p.font.color.rgb = self.DARK_GRAY
                p.word_wrap = True

                current_item_idx += 1

        if slide_number:
            self._add_footer(slide, slide_number)
        return slide

    # ========== TIMELINE SLIDES ==========
    
    def add_timeline_slide(self, title, image_url, milestones, slide_number=None):
        """Create timeline slide with image and milestones"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Header strip
        header_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(10), Inches(0.08)
        )
        header_line.fill.solid()
        header_line.fill.fore_color.rgb = self.TEAL
        header_line.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = self.FONT_NAME
        p.font.size = Pt(32)
        p.font.color.rgb = self.DARK_TEAL
        
        # Red line under title
        red_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5), Inches(1.6),
            Inches(4), Inches(0.05)
        )
        red_line.fill.solid()
        red_line.fill.fore_color.rgb = self.RED
        red_line.line.fill.background()

        # Image column
        img_left = Inches(0.5)
        img_top = Inches(2.2)
        img_width = Inches(3)
        img_height = Inches(4.5)

        try:
            slide.shapes.add_picture(image_url, img_left, img_top, width=img_width, height=img_height)
        except Exception as e:
            print(f"Could not load image {image_url}: {e}")
            placeholder = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                img_left, img_top,
                img_width, img_height
            )
            placeholder.fill.background()
            placeholder.line.fill.solid()
            placeholder.line.fill.fore_color.rgb = self.LIGHT_BLUE
            placeholder.line.width = Pt(1.5)
            text_frame = placeholder.text_frame
            text_frame.text = "Image Placeholder"
            text_frame.paragraphs[0].font.size = Pt(14)
            text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Milestones column
        milestones_left = img_left + img_width + Inches(0.5)
        milestones_width = Inches(9) - img_width - Inches(0.5) - Inches(0.5)
        milestones_top = Inches(2.2)

        milestone_count = min(len(milestones), 6)
        
        if milestone_count <= 3:
            event_font_size = Pt(18)
            spacing = Inches(0.4)
        elif milestone_count <= 4:
            event_font_size = Pt(16)
            spacing = Inches(0.3)
        else:
            event_font_size = Pt(14)
            spacing = Inches(0.25)
        
        current_y = milestones_top
        for i, milestone in enumerate(milestones[:6]):
            # Date box
            date_rect_width = Inches(1)
            date_rect_height = Inches(0.5)
            date_rect = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                milestones_left, current_y,
                date_rect_width, date_rect_height
            )
            date_rect.fill.solid()
            date_rect.fill.fore_color.rgb = self.TEAL
            date_rect.line.fill.background()

            date_text_box = slide.shapes.add_textbox(
                milestones_left, current_y,
                date_rect_width, date_rect_height
            )
            dtf = date_text_box.text_frame
            p = dtf.paragraphs[0]
            p.text = milestone['date']
            p.font.name = self.FONT_NAME
            p.font.size = Pt(10)
            p.font.color.rgb = self.WHITE
            p.alignment = PP_ALIGN.CENTER
            dtf.vertical_anchor = MSO_ANCHOR.MIDDLE

            # Event text
            event_text_box = slide.shapes.add_textbox(
                milestones_left + date_rect_width + Inches(0.2), current_y,
                milestones_width - date_rect_width - Inches(0.2), date_rect_height + Inches(0.2)
            )
            etf = event_text_box.text_frame
            p = etf.paragraphs[0]
            p.text = milestone['event']
            p.font.name = self.FONT_NAME
            p.font.size = event_font_size
            p.font.color.rgb = self.DARK_GRAY
            p.word_wrap = True

            current_y += (date_rect_height + spacing)
            
        if slide_number:
            self._add_footer(slide, slide_number)
        return slide

    # ========== COMPARISON SLIDES ==========
    
    def add_comparison_slide(self, title, left_side, right_side, slide_number=None, middle_side=None):
        """Create comparison slide with 2 or 3 columns"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Header strip
        header_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(10), Inches(0.08)
        )
        header_line.fill.solid()
        header_line.fill.fore_color.rgb = self.TEAL
        header_line.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = self.FONT_NAME
        p.font.size = Pt(32)
        p.font.color.rgb = self.DARK_TEAL
        
        # Red line under title
        red_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5), Inches(1.6),
            Inches(4), Inches(0.05)
        )
        red_line.fill.solid()
        red_line.fill.fore_color.rgb = self.RED
        red_line.line.fill.background()

        has_middle = middle_side is not None and middle_side.get('items')
        num_columns = 3 if has_middle else 2
        
        col_total_width = Inches(9)
        col_gap = Inches(0.5)
        col_effective_width = (col_total_width - (num_columns - 1) * col_gap) / num_columns
        
        # Determine font size
        all_items = left_side.get('items', []) + right_side.get('items', [])
        if has_middle:
            all_items += middle_side.get('items', [])
        
        max_item_count = min(max(len(all_items), 1), 6)
        font_size = Pt(16) if max_item_count <= 3 else Pt(14)
        bullet_spacing = Inches(0.15) if max_item_count <= 3 else Inches(0.1)

        def _add_comparison_column(slide, x, y, width, height, data, is_removable=False):
            # Column box
            col_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                x, y,
                width, height
            )
            col_shape.fill.background()
            col_shape.line.fill.solid()
            col_shape.line.fill.fore_color.rgb = self.LIGHT_BLUE
            col_shape.line.width = Pt(1.5)

            # Column Title
            title_box = slide.shapes.add_textbox(x + Inches(0.3), y + Inches(0.3), width - Inches(0.6), Inches(0.5))
            ttf = title_box.text_frame
            p = ttf.paragraphs[0]
            p.text = data['title']
            p.font.name = self.FONT_NAME
            p.font.size = Pt(20)
            p.font.color.rgb = self.DARK_TEAL
            p.word_wrap = True

            # Bullet points
            list_y = y + Inches(1)
            for item_text in data['items'][:6]:
                bullet_box = slide.shapes.add_textbox(x + Inches(0.5), list_y, width - Inches(1), Inches(0.5))
                btf = bullet_box.text_frame
                p = btf.paragraphs[0]
                p.text = f"• {item_text}"
                p.font.name = self.FONT_NAME
                p.font.size = font_size
                p.font.color.rgb = self.DARK_GRAY
                p.word_wrap = True
                list_y += (bullet_box.height + bullet_spacing)
            
            if is_removable:
                delete_cue = slide.shapes.add_textbox(
                    x + width - Inches(0.5), y + Inches(0.1),
                    Inches(0.4), Inches(0.4)
                )
                dtf = delete_cue.text_frame
                p = dtf.paragraphs[0]
                p.text = "X"
                p.font.name = self.FONT_NAME
                p.font.size = Pt(12)
                p.font.color.rgb = self.RED
                p.alignment = PP_ALIGN.CENTER
                dtf.vertical_anchor = MSO_ANCHOR.MIDDLE

        column_y_start = Inches(2.2)
        column_height = Inches(4.5)

        # Left column
        _add_comparison_column(
            slide,
            Inches(0.5), column_y_start,
            col_effective_width, column_height,
            left_side
        )

        if has_middle:
            # Middle column
            _add_comparison_column(
                slide,
                Inches(0.5) + col_effective_width + col_gap, column_y_start,
                col_effective_width, column_height,
                middle_side,
                is_removable=True
            )

        # Right column
        right_col_x = Inches(0.5) + (col_effective_width + col_gap) * (num_columns - 1)
        _add_comparison_column(
            slide,
            right_col_x, column_y_start,
            col_effective_width, column_height,
            right_side
        )
        
        if slide_number:
            self._add_footer(slide, slide_number)
        return slide

    # ========== PROCESS & WORKFLOW SLIDES ==========
    
    def add_process_steps_slide(self, title, steps, slide_number=None):
        """Create process steps slide with responsive grid and arrows"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Header strip
        header_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(10), Inches(0.08)
        )
        header_line.fill.solid()
        header_line.fill.fore_color.rgb = self.TEAL
        header_line.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = self.FONT_NAME
        p.font.size = Pt(32)
        p.font.color.rgb = self.DARK_TEAL
        
        # Red line under title
        red_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5), Inches(1.6),
            Inches(4), Inches(0.05)
        )
        red_line.fill.solid()
        red_line.fill.fore_color.rgb = self.RED
        red_line.line.fill.background()

        step_count = len(steps)
        if step_count == 0:
            if slide_number:
                self._add_footer(slide, slide_number)
            return slide

        # Layout
        max_steps_per_row = 4
        cols = min(step_count, 4)

        # Responsive sizing
        if step_count <= 3:
            box_width = Inches(2.2)
            box_height = Inches(2.2)
            label_font_size = Pt(20)
            desc_font_size = Pt(14)
            h_gap = Inches(0.5)
            v_gap = Inches(0.5)
        else:
            box_width = Inches(1.8)
            box_height = Inches(1.8)
            label_font_size = Pt(16)
            desc_font_size = Pt(12)
            h_gap = Inches(0.4)
            v_gap = Inches(0.4)

        total_width_available = Inches(9)
        
        if cols > 0:
            total_gap_width = (cols - 1) * h_gap
            actual_box_width = (total_width_available - total_gap_width) / cols
            box_width = min(box_width, actual_box_width)

        start_x = Inches(0.5)
        start_y = Inches(2.2)
        
        current_step_idx = 0
        while current_step_idx < step_count:
            row_start_y = start_y + (current_step_idx // cols) * (box_height + Inches(0.8) + v_gap)

            for c in range(cols):
                if current_step_idx >= step_count:
                    break
                
                step = steps[current_step_idx]
                x = start_x + c * (box_width + h_gap)
                y = row_start_y

                # Process box
                process_box = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    x, y,
                    box_width, box_height
                )
                process_box.fill.solid()
                process_box.fill.fore_color.rgb = self.TEAL
                process_box.line.fill.background()

                # Label
                label_text_box = slide.shapes.add_textbox(
                    x, y,
                    box_width, box_height
                )
                ltf = label_text_box.text_frame
                p = ltf.paragraphs[0]
                p.text = step['label']
                p.font.name = self.FONT_NAME
                p.font.size = label_font_size
                p.font.color.rgb = self.WHITE
                p.alignment = PP_ALIGN.CENTER
                ltf.vertical_anchor = MSO_ANCHOR.MIDDLE

                # Description
                desc_text_box = slide.shapes.add_textbox(
                    x, y + box_height + Inches(0.1),
                    box_width, Inches(0.7)
                )
                dtf = desc_text_box.text_frame
                p = dtf.paragraphs[0]
                p.text = step['description']
                p.font.name = self.FONT_NAME
                p.font.size = desc_font_size
                p.font.color.rgb = self.DARK_GRAY
                p.alignment = PP_ALIGN.CENTER
                p.word_wrap = True

                # Arrow
                if c < cols - 1 and current_step_idx < step_count - 1:
                    arrow = slide.shapes.add_shape(
                        MSO_SHAPE.RIGHT_ARROW,
                        x + box_width, y + box_height / 2 - Inches(0.05),
                        h_gap, Inches(0.1)
                    )
                    arrow.fill.solid()
                    arrow.fill.fore_color.rgb = self.RED
                    arrow.line.fill.background()
                
                current_step_idx += 1
            
            # Vertical arrow for next row
            if (current_step_idx % cols == 0 and current_step_idx < step_count):
                arrow_vertical = slide.shapes.add_shape(
                    MSO_SHAPE.DOWN_ARROW,
                    start_x + box_width / 2 - Inches(0.05), row_start_y + box_height + Inches(0.8) + Inches(0.1),
                    Inches(0.1), v_gap - Inches(0.2)
                )
                arrow_vertical.fill.solid()
                arrow_vertical.fill.fore_color.rgb = self.RED
                arrow_vertical.line.fill.background()

        if slide_number:
            self._add_footer(slide, slide_number)
        return slide

    # ========== TEAM SLIDES ==========
    
    def add_team_slide(self, title, members, slide_number=None):
        """Create team slide with member photos and info"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Header strip
        header_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(10), Inches(0.08)
        )
        header_line.fill.solid()
        header_line.fill.fore_color.rgb = self.TEAL
        header_line.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = self.FONT_NAME
        p.font.size = Pt(32)
        p.font.color.rgb = self.DARK_TEAL
        
        # Red line under title
        red_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5), Inches(1.6),
            Inches(4), Inches(0.05)
        )
        red_line.fill.solid()
        red_line.fill.fore_color.rgb = self.RED
        red_line.line.fill.background()

        member_count = len(members)
        if member_count == 0:
            if slide_number:
                self._add_footer(slide, slide_number)
            return slide

        # Grid configuration
        grid_configs = {
            1: (1, 1), 2: (1, 2), 3: (1, 3), 4: (1, 4),
            5: (2, 3), 6: (2, 3), 7: (2, 4), 8: (2, 4),
            9: (2, 5), 10: (2, 5)
        }
        num_rows, max_cols_per_row = grid_configs.get(member_count, (2, 5))

        # Sizing
        if member_count <= 4:
            img_size_inch = Inches(1.8)
            name_font_size = Pt(18)
            role_font_size = Pt(14)
            gap_h = Inches(0.6)
            gap_v = Inches(0.5)
        elif member_count <= 8:
            img_size_inch = Inches(1.5)
            name_font_size = Pt(16)
            role_font_size = Pt(12)
            gap_h = Inches(0.5)
            gap_v = Inches(0.4)
        else:
            img_size_inch = Inches(1.3)
            name_font_size = Pt(14)
            role_font_size = Pt(10)
            gap_h = Inches(0.4)
            gap_v = Inches(0.3)
        
        start_x = Inches(0.5)
        start_y = Inches(2.2)
        total_content_width = Inches(9)

        current_member_idx = 0
        for r in range(num_rows):
            cols_in_this_row = max_cols_per_row
            if member_count == 5 and r == 1:
                cols_in_this_row = 2
            elif member_count == 7 and r == 1:
                cols_in_this_row = 3
            elif member_count == 9 and r == 1:
                cols_in_this_row = 4
            elif member_count < max_cols_per_row:
                cols_in_this_row = member_count

            # Center items
            total_item_width_row = cols_in_this_row * img_size_inch.val + (cols_in_this_row - 1) * gap_h.val
            row_start_x = start_x.val + (total_content_width.val - total_item_width_row) / 2
            row_start_x = Inches(row_start_x)
            
            row_start_y = start_y + r * (img_size_inch + Inches(0.8) + gap_v)

            for c in range(cols_in_this_row):
                if current_member_idx >= member_count:
                    break
                
                member = members[current_member_idx]
                x = row_start_x + c * (img_size_inch + gap_h)
                y = row_start_y

                # Member image
                try:
                    slide.shapes.add_picture(member['image_url'], x, y, width=img_size_inch, height=img_size_inch)
                except Exception as e:
                    print(f"Could not load member image {member['image_url']}: {e}")
                    placeholder = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        x, y,
                        img_size_inch, img_size_inch
                    )
                    placeholder.fill.background()
                    placeholder.line.fill.solid()
                    placeholder.line.fill.fore_color.rgb = self.LIGHT_BLUE
                    placeholder.line.width = Pt(1.5)
                    text_frame = placeholder.text_frame
                    initials = "".join([n[0] for n in member['name'].split() if n])
                    text_frame.text = initials[:2].upper()
                    text_frame.paragraphs[0].font.size = Pt(name_font_size.val * 0.8)
                    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

                # Name
                name_box = slide.shapes.add_textbox(x, y + img_size_inch + Inches(0.1), img_size_inch, Inches(0.4))
                ntf = name_box.text_frame
                p = ntf.paragraphs[0]
                p.text = member['name']
                p.font.name = self.FONT_NAME
                p.font.size = name_font_size
                p.font.color.rgb = self.DARK_TEAL
                p.alignment = PP_ALIGN.CENTER
                p.word_wrap = True

                # Role
                role_box = slide.shapes.add_textbox(x, y + img_size_inch + Inches(0.1) + Inches(0.4), img_size_inch, Inches(0.3))
                rtf = role_box.text_frame
                p = rtf.paragraphs[0]
                p.text = member['role']
                p.font.name = self.FONT_NAME
                p.font.size = role_font_size
                p.font.color.rgb = self.DARK_GRAY
                p.alignment = PP_ALIGN.CENTER
                p.word_wrap = True

                current_member_idx += 1
        
        if slide_number:
            self._add_footer(slide, slide_number)
        return slide

    # ========== QUOTE SLIDES ==========
    
    def add_quote_slide(self, quote, author, role, slide_number=None):
        """Create inspirational quote slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Background
        bg_shape_blue = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(10), Inches(7.5)
        )
        bg_shape_blue.fill.solid()
        bg_shape_blue.fill.fore_color.rgb = self.TEAL
        bg_shape_blue.line.fill.background()

        # Red accent line
        red_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(10), Inches(0.08)
        )
        red_line.fill.solid()
        red_line.fill.fore_color.rgb = self.RED
        red_line.line.fill.background()

        # Quote text
        quote_len = len(quote)
        quote_font_size = Pt(36) if quote_len <= 100 else Pt(28) if quote_len <= 200 else Pt(22)
        
        quote_box = slide.shapes.add_textbox(
            Inches(1), Inches(2),
            Inches(8), Inches(3)
        )
        qtf = quote_box.text_frame
        p = qtf.paragraphs[0]
        p.text = f'"{quote}"'
        p.font.name = self.FONT_NAME
        p.font.size = quote_font_size
        p.font.color.rgb = self.WHITE
        p.alignment = PP_ALIGN.CENTER
        p.word_wrap = True
        qtf.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Author
        author_box = slide.shapes.add_textbox(
            Inches(1), Inches(5),
            Inches(8), Inches(0.5)
        )
        atf = author_box.text_frame
        p = atf.paragraphs[0]
        p.text = author
        p.font.name = self.FONT_NAME
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = self.WHITE
        p.alignment = PP_ALIGN.CENTER
        
        # Role
        role_box = slide.shapes.add_textbox(
            Inches(1), Inches(5.5),
            Inches(8), Inches(0.4)
        )
        rtf = role_box.text_frame
        p = rtf.paragraphs[0]
        p.text = role
        p.font.name = self.FONT_NAME
        p.font.size = Pt(14)
        p.font.color.rgb = self.WHITE
        p.alignment = PP_ALIGN.CENTER
        
        if slide_number:
            self._add_footer(slide, slide_number)
        return slide

    # ========== STATS SLIDES ==========
    
    def add_stats_slide(self, title, stats, slide_number=None):
        """Create statistics showcase slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Header strip
        header_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(10), Inches(0.08)
        )
        header_line.fill.solid()
        header_line.fill.fore_color.rgb = self.TEAL
        header_line.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = self.FONT_NAME
        p.font.size = Pt(32)
        p.font.color.rgb = self.DARK_TEAL
        
        # Red line under title
        red_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5), Inches(1.6),
            Inches(4), Inches(0.05)
        )
        red_line.fill.solid()
        red_line.fill.fore_color.rgb = self.RED
        red_line.line.fill.background()

        stat_count = len(stats)
        if stat_count == 0:
            if slide_number:
                self._add_footer(slide, slide_number)
            return slide

        # Grid config
        grid_config = {
            1: (1, 1), 2: (1, 2), 3: (1, 3),
            4: (2, 2), 5: (2, 3), 6: (2, 3)
        }
        num_rows, max_cols_per_row = grid_config.get(stat_count, (1, 3))
        
        # Sizing
        if stat_count <= 3:
            col_width = Inches(8.5 / stat_count)
            item_h_padding = Inches(0.8)
            number_font_size = Pt(48)
            label_font_size = Pt(18)
            gap = Inches(0.4)
        else:
            col_width = Inches(8.5 / 3)
            item_h_padding = Inches(0.5)
            number_font_size = Pt(36)
            label_font_size = Pt(14)
            gap = Inches(0.3)
        
        start_x = Inches(0.5)
        start_y = Inches(2.2)
        current_stat_idx = 0

        for r in range(num_rows):
            cols_in_this_row = max_cols_per_row
            if stat_count == 5 and r == 1:
                cols_in_this_row = 2
            elif stat_count <= 3:
                cols_in_this_row = stat_count

            row_start_y = start_y + r * (Inches(2.5) + gap)

            for c in range(cols_in_this_row):
                if current_stat_idx >= stat_count:
                    break

                stat = stats[current_stat_idx]
                
                x = start_x + c * (col_width + gap)
                y = row_start_y

                # Stat box
                stat_shape = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    x, y,
                    col_width, Inches(2.5)
                )
                stat_shape.fill.background()
                stat_shape.line.fill.solid()
                stat_shape.line.fill.fore_color.rgb = self.LIGHT_BLUE
                stat_shape.line.width = Pt(1.5)

                # Number
                number_text_box = slide.shapes.add_textbox(
                    x + Inches(0.1), y + Inches(0.5),
                    col_width - Inches(0.2), Inches(1)
                )
                ntf = number_text_box.text_frame
                p = ntf.paragraphs[0]
                p.text = stat['number']
                p.font.name = self.FONT_NAME
                p.font.size = number_font_size
                p.font.color.rgb = self.DARK_TEAL
                p.alignment = PP_ALIGN.CENTER
                
                # Label
                label_text_box = slide.shapes.add_textbox(
                    x + Inches(0.1), y + Inches(1.5),
                    col_width - Inches(0.2), Inches(0.7)
                )
                ltf = label_text_box.text_frame
                p = ltf.paragraphs[0]
                p.text = stat['label']
                p.font.name = self.FONT_NAME
                p.font.size = label_font_size
                p.font.color.rgb = self.DARK_GRAY
                p.alignment = PP_ALIGN.CENTER
                p.word_wrap = True

                current_stat_idx += 1

        if slide_number:
            self._add_footer(slide, slide_number)
        return slide

    # ========== CONTACT INFO SLIDES ==========
    
    def add_contact_info_slide(self, title, image_url, contact_details, slide_number=None):
        """Create contact information slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Header strip
        header_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(10), Inches(0.08)
        )
        header_line.fill.solid()
        header_line.fill.fore_color.rgb = self.TEAL
        header_line.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = self.FONT_NAME
        p.font.size = Pt(32)
        p.font.color.rgb = self.DARK_TEAL
        
        # Red line under title
        red_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5), Inches(1.6),
            Inches(4), Inches(0.05)
        )
        red_line.fill.solid()
        red_line.fill.fore_color.rgb = self.RED
        red_line.line.fill.background()

        # Image
        img_left = Inches(0.5)
        img_top = Inches(2.2)
        img_width = Inches(4.5)
        img_height = Inches(4.5)
        
        try:
            slide.shapes.add_picture(image_url, img_left, img_top, width=img_width, height=img_height)
        except Exception as e:
            print(f"Could not load image {image_url}: {e}")
            placeholder = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                img_left, img_top,
                img_width, img_height
            )
            placeholder.fill.background()
            placeholder.line.fill.solid()
            placeholder.line.fill.fore_color.rgb = self.LIGHT_BLUE
            placeholder.line.width = Pt(1.5)
            text_frame = placeholder.text_frame
            text_frame.text = "Contact Image"
            text_frame.paragraphs[0].font.size = Pt(14)
            text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Contact details
        details_left = Inches(0.5) + img_width + Inches(0.5)
        details_top = Inches(2.2)
        details_width = Inches(10) - details_left - Inches(0.5)

        detail_count = len(contact_details)
        detail_spacing = Inches(0.3) if detail_count <= 2 else Inches(0.2)
        
        current_y = details_top
        for detail in contact_details:
            # Icon box
            icon_box_width = Inches(0.6)
            icon_box_height = Inches(0.6)
            icon_rect = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                details_left, current_y,
                icon_box_width, icon_box_height
            )
            icon_rect.fill.solid()
            icon_rect.fill.fore_color.rgb = self.RED
            icon_rect.line.fill.background()

            icon_text_box = slide.shapes.add_textbox(
                details_left, current_y,
                icon_box_width, icon_box_height
            )
            itf = icon_text_box.text_frame
            itf.text = detail.get('icon', '★')
            itf.paragraphs[0].font.name = self.FONT_NAME
            itf.paragraphs[0].font.size = Pt(14)
            itf.paragraphs[0].font.color.rgb = self.WHITE
            itf.paragraphs[0].alignment = PP_ALIGN.CENTER
            itf.vertical_anchor = MSO_ANCHOR.MIDDLE

            # Label and Value
            text_x = details_left + icon_box_width + Inches(0.2)
            text_width = details_width - icon_box_width - Inches(0.2)

            label_box = slide.shapes.add_textbox(text_x, current_y, text_width, Inches(0.3))
            ltf = label_box.text_frame
            p = ltf.paragraphs[0]
            p.text = detail['label'].upper()
            p.font.name = self.FONT_NAME
            p.font.size = Pt(9)
            p.font.color.rgb = self.DARK_GRAY
            
            value_box = slide.shapes.add_textbox(text_x, current_y + Inches(0.3), text_width, Inches(0.5))
            vtf = value_box.text_frame
            p = vtf.paragraphs[0]
            p.text = detail['value']
            p.font.name = self.FONT_NAME
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = self.DARK_TEAL
            
            current_y += (icon_box_height + detail_spacing)
            
        if slide_number:
            self._add_footer(slide, slide_number)
        return slide

    # ========== IMAGE TEXT SPLIT SLIDES ==========
    
    def add_image_text_split_slide(self, title, image_url, content, image_position='left', slide_number=None):
        """Create slide with image and text side by side"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Header strip
        header_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(10), Inches(0.08)
        )
        header_line.fill.solid()
        header_line.fill.fore_color.rgb = self.TEAL
        header_line.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = self.FONT_NAME
        p.font.size = Pt(32)
        p.font.color.rgb = self.DARK_TEAL
        
        # Red line under title
        red_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5), Inches(1.6),
            Inches(4), Inches(0.05)
        )
        red_line.fill.solid()
        red_line.fill.fore_color.rgb = self.RED
        red_line.line.fill.background()

        # Layout: 1/3 for image, 2/3 for text
        img_width = Inches(3)
        text_width = Inches(6)
        
        main_content_top = Inches(2.2)
        main_content_height = Inches(4.5)

        image_left_x = Inches(0.5) if image_position == 'left' else Inches(0.5) + text_width + Inches(0.5)
        text_left_x = Inches(0.5) if image_position == 'right' else Inches(0.5) + img_width + Inches(0.5)

        # Image
        try:
            slide.shapes.add_picture(image_url, image_left_x, main_content_top, width=img_width, height=main_content_height)
        except Exception as e:
            print(f"Could not load image {image_url}: {e}")
            placeholder = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                image_left_x, main_content_top,
                img_width, main_content_height
            )
            placeholder.fill.background()
            placeholder.line.fill.solid()
            placeholder.line.fill.fore_color.rgb = self.LIGHT_BLUE
            placeholder.line.width = Pt(1.5)
            text_frame = placeholder.text_frame
            text_frame.text = "Image Placeholder"
            text_frame.paragraphs[0].font.size = Pt(14)
            text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Text content
        heading_text = content.get('heading', '')
        paragraphs = content.get('paragraphs', [])

        para_count = len(paragraphs)
        para_font_size = Pt(18) if para_count <= 2 else Pt(16) if para_count <= 3 else Pt(14)
        para_spacing = Inches(0.2) if para_count <= 2 else Inches(0.15) if para_count <= 3 else Inches(0.1)

        current_y = main_content_top

        if heading_text:
            heading_box = slide.shapes.add_textbox(text_left_x, current_y, text_width, Inches(0.5))
            htf = heading_box.text_frame
            p = htf.paragraphs[0]
            p.text = heading_text
            p.font.name = self.FONT_NAME
            p.font.size = Pt(22)
            p.font.bold = True
            p.font.color.rgb = self.DARK_TEAL
            current_y += Inches(0.7)

        for para in paragraphs:
            text_box = slide.shapes.add_textbox(text_left_x, current_y, text_width, Inches(1))
            tf = text_box.text_frame
            p = tf.paragraphs[0]
            p.text = para
            p.font.name = self.FONT_NAME
            p.font.size = para_font_size
            p.font.color.rgb = self.DARK_GRAY
            p.word_wrap = True
            current_y += (text_box.height + para_spacing)
        
        if slide_number:
            self._add_footer(slide, slide_number)
        return slide

    # ========== CHART SLIDES ==========
    
    def add_chart_slide(self, title, chart_type, chart_data, slide_number=None):
        """Create chart slide (line or bar)"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Header strip
        header_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(10), Inches(0.08)
        )
        header_line.fill.solid()
        header_line.fill.fore_color.rgb = self.TEAL
        header_line.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = self.FONT_NAME
        p.font.size = Pt(32)
        p.font.color.rgb = self.DARK_TEAL
        
        # Red line under title
        red_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5), Inches(1.6),
            Inches(4), Inches(0.05)
        )
        red_line.fill.solid()
        red_line.fill.fore_color.rgb = self.RED
        red_line.line.fill.background()

        # Chart area
        chart_x, chart_y, chart_width, chart_height = Inches(0.5), Inches(2.2), Inches(9), Inches(4.5)

        category_names = chart_data['categories']
        series_data = []
        for series in chart_data['series']:
            series_data.append((series['name'], series['values']))

        graphic_frame = None
        if chart_type == 'line':
            chart_type_enum = XL_CHART_TYPE.LINE
        elif chart_type == 'bar':
            chart_type_enum = XL_CHART_TYPE.COLUMN_CLUSTERED
        else:
            print(f"Unsupported chart type: {chart_type}")
            if slide_number:
                self._add_footer(slide, slide_number)
            return slide

        chart_data_obj = CategoryChartData()
        chart_data_obj.categories = category_names
        for name, values in series_data:
            chart_data_obj.add_series(name, values)
        
        graphic_frame = slide.shapes.add_chart(
            chart_type_enum, chart_x, chart_y, chart_width, chart_height, chart_data_obj
        )

        # Chart styling
        if graphic_frame:
            chart = graphic_frame.chart
            chart.has_legend = True
            chart.legend.include_in_layout = False
            chart.legend.position = 1
            
            category_axis = chart.category_axis
            category_axis.has_major_gridlines = False
            category_axis.tick_labels.font.name = self.FONT_NAME
            category_axis.tick_labels.font.size = Pt(10)
            category_axis.tick_labels.font.color.rgb = self.DARK_GRAY

            value_axis = chart.value_axis
            value_axis.has_major_gridlines = True
            value_axis.major_gridlines.format.line.color.rgb = self.LIGHT_GRAY
            value_axis.tick_labels.font.name = self.FONT_NAME
            value_axis.tick_labels.font.size = Pt(10)
            value_axis.tick_labels.font.color.rgb = self.DARK_GRAY

            # Series colors
            for i, series in enumerate(chart.series):
                fill = series.format.fill
                fill.solid()
                if i == 0:
                    fill.fore_color.rgb = self.TEAL
                else:
                    fill.fore_color.rgb = self.LIGHT_BLUE

        if slide_number:
            self._add_footer(slide, slide_number)
        return slide

    # ========== SAVE METHOD ==========
    
    def save(self, path):
        """Save presentation to file"""
        self.prs.save(path)
        print(f"✅ Presentation saved: {path}")


# Example Usage
if __name__ == '__main__':
    presentation = CorporatePresentation()

    # Title Slide
    presentation.add_title_slide(
        title="Your Presentation Title",
        subtitle="A captivating subtitle for your audience",
        slide_number=1
    )

    # Table of Contents
    presentation.add_table_of_contents(
        sections=[
            "Executive Summary",
            "Market Analysis",
            "Our Solution",
            "Product Features",
            "Timeline & Roadmap",
            "Team Introduction",
            "Key Metrics",
            "Contact Information"
        ],
        slide_number=2
    )

    # Content with Icons
    presentation.add_content_with_icons_slide(
        title="Key Features & Benefits",
        items=[
            {'icon': '💡', 'text': 'Innovative Solutions for Modern Problems'},
            {'icon': '🚀', 'text': 'Accelerate Your Business Growth and Efficiency'},
            {'icon': '👥', 'text': 'Dedicated Support and Expert Team Collaboration'},
            {'icon': '✓', 'text': 'Proven Track Record of Success and Reliability'},
        ],
        slide_number=3
    )

    # Split Slide
    presentation.add_split_slide(
        title="Understanding Our Approach",
        paragraphs=[
            "Our strategic framework is built on a foundation of rigorous research and adaptive methodologies.",
            "We prioritize collaborative development, integrating client insights at every stage.",
            "Continuous improvement is at the core of our operations."
        ],
        slide_number=4
    )

    # Market Opportunities
    presentation.add_market_opportunities_slide(
        title="Untapped Market Opportunities",
        items=[
            {'icon': '🌍', 'title': 'Global Expansion', 'text': 'Tap into emerging international markets with high growth potential.'},
            {'icon': '📱', 'title': 'Mobile Integration', 'text': 'Develop mobile-first solutions to capture the growing smartphone user base.'},
            {'icon': '⚡', 'title': 'AI Automation', 'text': 'Leverage AI to automate processes and enhance decision-making.'},
            {'icon': '🛡️', 'title': 'Enhanced Security', 'text': 'Offer robust cybersecurity features to meet increasing demands.'},
        ],
        slide_number=5
    )

    # Stats Slide
    presentation.add_stats_slide(
        title="Key Performance Indicators",
        stats=[
            {'number': '99.9%', 'label': 'Uptime Reliability'},
            {'number': '24/7', 'label': 'Global Support'},
            {'number': '10M+', 'label': 'Active Users'},
        ],
        slide_number=6
    )

    # Quote Slide
    presentation.add_quote_slide(
        quote="The only way to do great work is to love what you do.",
        author="Steve Jobs",
        role="Co-founder of Apple Inc.",
        slide_number=7
    )

    presentation.save("/mnt/user-data/outputs/corporate_presentation_enhanced.pptx")


